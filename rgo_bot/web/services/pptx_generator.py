"""Presentation generation service.

Voice → Whisper → Claude (plan + content) → python-pptx → send PPTX via bot.
"""
from __future__ import annotations

import io
import json

from aiogram import Bot
from aiogram.types import BufferedInputFile
from loguru import logger

from rgo_bot.bot.config import settings
from rgo_bot.db.base import async_session
from rgo_bot.web.app import update_task


async def generate_presentation(
    audio_data: bytes,
    filename: str,
    bot: Bot,
    user_id: int,
    task_id: str,
) -> dict:
    """Full pipeline: voice → transcript → plan → slides → PPTX.

    Args:
        audio_data: raw audio bytes with presentation requirements
        filename: original filename
        bot: aiogram Bot instance
        user_id: Telegram user ID
        task_id: task ID for progress updates

    Returns:
        dict with status
    """
    # Step 1: Transcribe
    update_task(task_id, step="transcribing")
    from rgo_bot.web.services.meeting_summarizer import _transcribe_bytes

    transcript = await _transcribe_bytes(audio_data, filename)
    if not transcript:
        await bot.send_message(
            settings.admin_telegram_id,
            "❌ Не удалось расшифровать аудио для презентации",
        )
        return {"error": "Transcription failed"}

    # Step 2: Load user preferences
    from rgo_bot.db.crud.presentation_preferences import get_preferences

    async with async_session() as session:
        prefs_record = await get_preferences(session, user_id)

    prefs_json = prefs_record.preferences_json if prefs_record else {}

    # Step 3: Generate presentation plan via Claude
    update_task(task_id, step="planning")
    plan = await _generate_plan(transcript, prefs_json)
    if not plan:
        await bot.send_message(
            settings.admin_telegram_id,
            "❌ Ошибка планирования презентации",
        )
        return {"error": "Planning failed"}

    # Step 4: Generate slide content
    update_task(task_id, step="generating")
    slides = await _generate_slide_content(plan)
    if not slides:
        await bot.send_message(
            settings.admin_telegram_id,
            "❌ Ошибка генерации контента слайдов",
        )
        return {"error": "Slide content generation failed"}

    # Step 5: Build PPTX
    update_task(task_id, step="building_pptx")
    pptx_bytes = _build_pptx(plan, slides)

    # Step 6: Send via bot
    title = plan.get("title", "Презентация")
    safe_title = "".join(c for c in title if c.isalnum() or c in " _-").strip()[:50]
    doc = BufferedInputFile(
        file=pptx_bytes,
        filename=f"{safe_title or 'Презентация'}.pptx",
    )
    await bot.send_document(
        settings.admin_telegram_id,
        document=doc,
        caption=f"📑 <b>{title}</b>\n\nСлайдов: {len(slides)}",
    )

    # Step 7: Save updated preferences
    updated_prefs = plan.get("updated_preferences", {})
    if updated_prefs:
        from rgo_bot.db.crud.presentation_preferences import upsert_preferences

        async with async_session() as session:
            await upsert_preferences(session, user_id, updated_prefs)

    return {"status": "ok", "title": title, "slides": len(slides)}


async def _generate_plan(transcript: str, preferences: dict) -> dict | None:
    """Generate presentation plan via Claude."""
    try:
        from rgo_bot.bot.services.claude_client import claude_client, load_prompt

        system_prompt = load_prompt("system")
        user_prompt = load_prompt("presentation_plan").format(
            user_request=transcript,
            preferences=json.dumps(preferences, ensure_ascii=False) if preferences else "Нет",
        )

        response = await claude_client.complete(
            system_prompt=system_prompt,
            user_prompt=user_prompt,
            max_tokens=2048,
            temperature=0.4,
            call_type="presentation_plan",
        )

        # Parse JSON from response
        text = response.text.strip()
        # Strip markdown code fences if present
        if text.startswith("```"):
            text = text.split("\n", 1)[1] if "\n" in text else text[3:]
            if text.endswith("```"):
                text = text[:-3]
            text = text.strip()

        return json.loads(text)

    except (json.JSONDecodeError, Exception):
        logger.exception("preza_plan_error")
        return None


async def _generate_slide_content(plan: dict) -> list[dict] | None:
    """Generate detailed slide content via Claude."""
    try:
        from rgo_bot.bot.services.claude_client import claude_client, load_prompt

        system_prompt = load_prompt("system")
        user_prompt = load_prompt("presentation_slides").format(
            plan_json=json.dumps(plan, ensure_ascii=False, indent=2)
        )

        response = await claude_client.complete(
            system_prompt=system_prompt,
            user_prompt=user_prompt,
            max_tokens=4096,
            temperature=0.3,
            call_type="presentation_slides",
        )

        text = response.text.strip()
        if text.startswith("```"):
            text = text.split("\n", 1)[1] if "\n" in text else text[3:]
            if text.endswith("```"):
                text = text[:-3]
            text = text.strip()

        return json.loads(text)

    except (json.JSONDecodeError, Exception):
        logger.exception("preza_slides_error")
        return None


def _build_pptx(plan: dict, slides: list[dict]) -> bytes:
    """Build PPTX in Sber corporate green style.

    Colors: bg=F0F2EF, dark=1A2B1E, accent=21A038, text=1A1A1A, gray=5C5C5C
    Font: Calibri everywhere
    Structure: first/last slide dark (1A2B1E), content slides light (F0F2EF)
    Footer: green line + СБЕР left + title+page right
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Corporate color palette
    CLR_BG = RGBColor(0xF0, 0xF2, 0xEF)       # Light background
    CLR_DARK = RGBColor(0x1A, 0x2B, 0x1E)      # Title/final slide bg
    CLR_ACCENT = RGBColor(0x21, 0xA0, 0x38)     # Green accent
    CLR_TEXT = RGBColor(0x1A, 0x1A, 0x1A)        # Main text
    CLR_GRAY = RGBColor(0x5C, 0x5C, 0x5C)       # Secondary text
    CLR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)       # White
    CLR_FOOTER_GRAY = RGBColor(0x9A, 0x9A, 0x9A) # Footer page number
    FONT = "Calibri"

    slide_layout = prs.slide_layouts[6]  # Blank
    title_text = plan.get("title", "Презентация")
    subtitle_text = plan.get("subtitle", "")
    total_slides = len(slides) + 2  # +title +final

    def _set_font(paragraph, size, color, bold=False, font_name=FONT):
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = color
        paragraph.font.bold = bold
        paragraph.font.name = font_name

    def _add_footer(slide, slide_num):
        """Add corporate footer: green line + СБЕР + page number."""
        # Green line
        line = slide.shapes.add_shape(
            1, Inches(0.8), Inches(6.85),
            Inches(11.733), Inches(0.02),
        )
        line.fill.solid()
        line.fill.fore_color.rgb = CLR_ACCENT
        line.line.fill.background()

        # "СБЕР" left
        txBox = slide.shapes.add_textbox(
            Inches(0.8), Inches(6.95), Inches(2), Inches(0.3)
        )
        p = txBox.text_frame.paragraphs[0]
        p.text = "СБЕР"
        _set_font(p, 13, CLR_ACCENT, bold=True)

        # Title + page number right
        txBox = slide.shapes.add_textbox(
            Inches(9), Inches(6.95), Inches(3.533), Inches(0.3)
        )
        p = txBox.text_frame.paragraphs[0]
        p.text = f"{title_text}  |  {slide_num}/{total_slides}"
        _set_font(p, 9, CLR_FOOTER_GRAY)
        p.alignment = PP_ALIGN.RIGHT

    # ── TITLE SLIDE (dark background) ──────────────────

    slide = prs.slides.add_slide(slide_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = CLR_DARK

    # Title
    txBox = slide.shapes.add_textbox(
        Inches(1), Inches(2.2), Inches(11), Inches(1.5)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    _set_font(p, 40, CLR_WHITE, bold=True)
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        _set_font(p2, 20, CLR_ACCENT)
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(16)

    # Accent line under title
    line = slide.shapes.add_shape(
        1, Inches(5.5), Inches(3.9), Inches(2.333), Inches(0.04),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = CLR_ACCENT
    line.line.fill.background()

    _add_footer(slide, 1)

    # ── CONTENT SLIDES (light background) ──────────────

    for idx, slide_data in enumerate(slides):
        slide_num = idx + 2
        slide = prs.slides.add_slide(slide_layout)

        # Light background
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = CLR_BG

        # Green accent bar at top
        bar = slide.shapes.add_shape(
            1, Inches(0), Inches(0),
            Inches(0.08), Inches(7.5),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = CLR_ACCENT
        bar.line.fill.background()

        # Slide title
        stitle = slide_data.get("title", "")
        txBox = slide.shapes.add_textbox(
            Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.8)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = stitle
        _set_font(p, 36, CLR_TEXT, bold=True)

        # Content area
        body_text = slide_data.get("body_text", "")
        bullets = slide_data.get("bullet_points", [])
        metrics = slide_data.get("metrics", [])  # e.g. [{"value": "+25%", "label": "рост"}]

        content_top = Inches(1.5)
        txBox = slide.shapes.add_textbox(
            Inches(0.8), content_top, Inches(11.5), Inches(4.8)
        )
        tf = txBox.text_frame
        tf.word_wrap = True

        if body_text:
            p = tf.paragraphs[0]
            p.text = body_text
            _set_font(p, 14, CLR_GRAY)
            p.space_after = Pt(12)

        for i, bullet in enumerate(bullets):
            p = tf.add_paragraph() if (body_text or i > 0) else tf.paragraphs[0]
            p.text = f"•  {bullet}"
            _set_font(p, 14, CLR_GRAY)
            p.space_before = Pt(8)
            p.space_after = Pt(4)

        # Metrics (big numbers) if provided by Claude
        if metrics:
            metric_left = Inches(0.8)
            for m_idx, metric in enumerate(metrics):
                mx = metric_left + Inches(m_idx * 3.5)
                # Value
                txBox = slide.shapes.add_textbox(
                    mx, Inches(4.2), Inches(3), Inches(0.8)
                )
                p = txBox.text_frame.paragraphs[0]
                p.text = str(metric.get("value", ""))
                _set_font(p, 56, CLR_TEXT, bold=True)
                # Label
                txBox = slide.shapes.add_textbox(
                    mx, Inches(5.0), Inches(3), Inches(0.3)
                )
                p = txBox.text_frame.paragraphs[0]
                p.text = str(metric.get("label", ""))
                _set_font(p, 12, CLR_GRAY)

        # Badge (dynamic indicator) if provided
        badge = slide_data.get("badge", "")  # e.g. "+15%"
        if badge:
            shape = slide.shapes.add_shape(
                5,  # Rounded rectangle
                Inches(10.5), Inches(0.45), Inches(1.8), Inches(0.4),
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = CLR_ACCENT
            shape.line.fill.background()
            tf = shape.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.text = badge
            _set_font(p, 11, CLR_WHITE, bold=True)
            p.alignment = PP_ALIGN.CENTER

        # Speaker notes
        notes = slide_data.get("speaker_notes", "")
        if notes:
            notes_slide = slide.notes_slide
            notes_tf = notes_slide.notes_text_frame
            notes_tf.text = notes

        _add_footer(slide, slide_num)

    # ── FINAL SLIDE (dark background) ──────────────────

    slide = prs.slides.add_slide(slide_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = CLR_DARK

    txBox = slide.shapes.add_textbox(
        Inches(1), Inches(2.8), Inches(11), Inches(1.2)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Спасибо за внимание"
    _set_font(p, 40, CLR_WHITE, bold=True)
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = title_text
    _set_font(p2, 18, CLR_ACCENT)
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(16)

    # Accent line
    line = slide.shapes.add_shape(
        1, Inches(5.5), Inches(4.2), Inches(2.333), Inches(0.04),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = CLR_ACCENT
    line.line.fill.background()

    _add_footer(slide, total_slides)

    # Save to bytes
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
